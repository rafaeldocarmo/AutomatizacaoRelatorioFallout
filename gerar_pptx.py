# -*- coding: utf-8 -*-
"""Blocos de montagem dos slides em PowerPoint editável (tabelas nativas +
gráfico como imagem), consumidos por gerar_slide3col.gerar_pptx_completo().
"""
import os, tempfile
import matplotlib
matplotlib.use("Agg")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

import fallout_core

RED = RGBColor(0xC0, 0x39, 0x2B); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xF2, 0xF2, 0xF2); GRAY2 = RGBColor(0xF5, 0xF5, 0xF5)
FG = RGBColor(0x22, 0x22, 0x22)


# ── Extração dos dados via pipeline ──────────────────────────────────────────
def _montar_dados(jornada, ns):
    """Monta (data, caminho_chart_png) a partir de um resultado já calculado
    por pipeline.gerar_relatorio()."""
    import pandas as pd
    resumo = ns["resumo"]; ML = ns["MESES_LABEL"]; reducao = ns["reducao"]
    meses_vol = ns["meses_vol"]
    fi = lambda v: f"{int(v):,}".replace(",", ".")
    fp = lambda v: f"{v:.2f}%".replace(".", ",")

    volume = {"months": [ML[m].capitalize() for m in meses_vol], "rows": [
        ["Vendas com Sucesso"] + [fi(resumo.loc[m, "Sucessos"]) for m in meses_vol],
        ["Falha (Análise Técnica)"] + [fi(resumo.loc[m, "Falhas"]) for m in meses_vol],
        ["Fallout Rate"] + [fp(resumo.loc[m, "Pct"]) for m in meses_vol],
    ]}
    dist = {"title": f"Distribuição Fallout ({ns['fallout_pct']:.2f}%)", "rows": [
        {"label": lbl, "val": fp(val), "sub": sub, "bold": bold}
        for (lbl, val, sub, bold) in ns["dist_items"]]}
    plan_rows = []
    for _, r in reducao.iterrows():
        dfts = r["DFTs"].split("\n")
        plan_rows.append({"date": r["MilestoneDate"].strftime("%d/%m/%Y"),
                          "dfts": [{"id": d.strip(), "pct": fp(r["Pct"] / len(dfts))} for d in dfts]})
    plan = {"title": f"Planejamento Redução ({reducao['Pct_plena'].sum():.2f}% pleno)", "rows": plan_rows}

    COL_NAMES = ns["COL_NAMES"]; meses_tab = ns["meses_tab"]
    pivot2 = ns["pivot2"].reset_index(drop=True)
    fmt_id = ns["fmt_id"]; fmt_ms = ns["fmt_milestone"]; pct_m = ns["pct_m"]
    n_fixas = len(COL_NAMES) - len(meses_tab)
    det_rows = []
    suc = ["Sucesso"] + [""] * (n_fixas - 1)
    for m in meses_tab:
        sv = pct_m(resumo.loc[m, "Sucessos"] if m in resumo.index else 0, m)
        suc.append(f"{sv:.2f}%".replace(".", ","))
    det_rows.append({"cells": suc, "kind": "sucesso"})
    for ri in range(len(pivot2)):
        row = pivot2.iloc[ri]
        tipo = str(row["DFT_Type"]).strip() if pd.notna(row["DFT_Type"]) else ""
        cells = [
            "Falha" if ri == 0 else "",
            "Em Tratamento/Avaliação pela Squad" if ri == 0 else "",
            "US" if tipo == "User Story" else "Defeito",
            fmt_id({"DefectKey": row["DefectKey"], "DFT_BugfixMilestone": row["DFT_BugfixMilestone"]}),
            fmt_ms(row["DFT_BugfixMilestone"]),
            str(row["DFT_Name"]) if pd.notna(row["DFT_Name"]) else "",
            str(row["DFT_Team"])[:35] if pd.notna(row["DFT_Team"]) else "",
            str(row["DFT_Phase"])[:25] if pd.notna(row["DFT_Phase"]) else "",
        ]
        for m in meses_tab:
            val = int(row.loc[m]) if m in pivot2.columns else 0
            cells.append(f"{pct_m(val, m):.2f}%".replace(".", ","))
        det_rows.append({"cells": cells, "kind": "dft"})
    detalhe = {"cols": COL_NAMES, "rows": det_rows}

    # recorta o gráfico
    fig = ns["fig"]; axc = ns["ax_chart"]
    fig.canvas.draw()
    ext = axc.get_tightbbox(fig.canvas.get_renderer()).transformed(fig.dpi_scale_trans.inverted())
    ext = ext.expanded(1.02, 1.0)
    chart_png = os.path.join(tempfile.mkdtemp(), "chart.png")
    fig.savefig(chart_png, dpi=200, bbox_inches=ext, pad_inches=0.08, facecolor="white")

    data = {"title": f"{fallout_core.nome_exibicao(jornada)} – Realização e Projeção de Fallout",
            "volume": volume, "dist": dist, "plan": plan, "detalhe": detalhe}
    return data, chart_png


# ── Montagem do slide ────────────────────────────────────────────────────────
def _no_style(tbl):
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0"); tblPr.set("bandRow", "0")
    for sid in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(sid)

def _cell(cell, text, size, bold=False, fg=FG, bg=WHITE, align="center", valign="middle"):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)
    cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE if valign == "middle" else MSO_ANCHOR.TOP
    tf = cell.text_frame; tf.word_wrap = True
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = fg; r.font.name = "Calibri"

def _no_border(cell):
    """Remove as bordas da célula (deixa como o PNG, sem grade)."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    for tag in ("lnB", "lnT", "lnR", "lnL"):   # insere no início → ordem final lnL,lnR,lnT,lnB
        ln = parse_xml('<a:%s %s w="0" cap="flat"><a:noFill/></a:%s>' % (tag, nsdecls("a"), tag))
        tcPr.insert(0, ln)

def _sem_bordas(tbl):
    for row in tbl.rows:
        for cell in row.cells:
            _no_border(cell)

def _borda_clara(tbl, hexcor="D9D9D9", w=6350):
    """Bordas finas em cinza claro (como o PNG) em todas as células."""
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                for el in tcPr.findall(qn(tag)):
                    tcPr.remove(el)
            for tag in ("lnB", "lnT", "lnR", "lnL"):
                ln = parse_xml(
                    '<a:%s %s w="%d" cap="flat"><a:solidFill><a:srgbClr val="%s"/></a:solidFill></a:%s>'
                    % (tag, nsdecls("a"), w, hexcor, tag))
                tcPr.insert(0, ln)

def _widths(tbl, ws):
    for i, w in enumerate(ws): tbl.columns[i].width = Inches(w)
def _rowh(tbl, h):
    for r in tbl.rows: r.height = Inches(h)

def _montar_slide(prs, d, chart_png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.55))
    tb.fill.solid(); tb.fill.fore_color.rgb = RED; tb.line.fill.background()
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = d["title"]
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

    def add_tbl(nr, nc, x, y, w, h):
        gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
        _no_style(gf.table); return gf.table

    LX, LW = 0.18, 5.75; y = 0.62
    vm = d["volume"]["months"]; drows_d = d["dist"]["rows"]; prows = d["plan"]["rows"]
    n_plan = 1 + sum(1 + len(r["dfts"]) for r in prows)
    left_total = (2 + len(d["volume"]["rows"])) + (1 + len(drows_d)) + n_plan
    GAP = 0.04
    # Fonte fixa em 8pt. A linha fica no mínimo que 8pt exige em vez de esticar
    # para ocupar o espaço: o bloco da esquerda disputa altura com a tabela de
    # detalhe lá embaixo, e esticar aqui comia as linhas de defeito de lá.
    FD, FT = 8.0, 9.0
    RH = 0.16

    vt = add_tbl(2 + len(d["volume"]["rows"]), 1 + len(vm), LX, y, LW, RH * (2 + len(d["volume"]["rows"])))
    vlabw = 2.5; vcolw = (LW - vlabw) / len(vm)
    _widths(vt, [vlabw] + [vcolw] * len(vm)); _rowh(vt, RH)
    _cell(vt.cell(0, 0), "", FT, bg=RED)
    vt.cell(0, 1).merge(vt.cell(0, len(vm)))
    _cell(vt.cell(0, 1), "Volume de Pedidos", FT, bold=True, fg=WHITE, bg=RED)
    _cell(vt.cell(1, 0), "", FT, bg=RED)
    for i, ml in enumerate(vm): _cell(vt.cell(1, 1 + i), ml, FT, bold=True, fg=WHITE, bg=RED)
    for ri, row in enumerate(d["volume"]["rows"]):
        bg = WHITE if ri % 2 == 0 else GRAY
        _cell(vt.cell(2 + ri, 0), row[0], FD, bg=bg, align="left")
        for i in range(len(vm)): _cell(vt.cell(2 + ri, 1 + i), row[1 + i], FD, bg=bg)
    y += RH * (2 + len(d["volume"]["rows"])) + GAP

    dt = add_tbl(1 + len(drows_d), 2, LX, y, LW, RH * (1 + len(drows_d)))
    _widths(dt, [LW - 1.5, 1.5]); _rowh(dt, RH)
    dt.cell(0, 0).merge(dt.cell(0, 1))
    _cell(dt.cell(0, 0), d["dist"]["title"], FT, bold=True, fg=WHITE, bg=RED)
    tgl = False
    for ri, row in enumerate(drows_d):
        bg = WHITE if not tgl else GRAY; tgl = not tgl
        lbl = ("     " + row["label"]) if row["sub"] else row["label"]
        _cell(dt.cell(1 + ri, 0), lbl, FD, bold=row["bold"], bg=bg, align="left")
        _cell(dt.cell(1 + ri, 1), row["val"], FD, bold=row["bold"], bg=bg, align="right")
    y += RH * (1 + len(drows_d)) + GAP

    pt = add_tbl(n_plan, 2, LX, y, LW, RH * n_plan)
    _widths(pt, [LW - 1.5, 1.5]); _rowh(pt, RH)
    pt.cell(0, 0).merge(pt.cell(0, 1))
    _cell(pt.cell(0, 0), d["plan"]["title"], FT, bold=True, fg=WHITE, bg=RED)
    ridx = 1
    for pr in prows:
        pt.cell(ridx, 0).merge(pt.cell(ridx, 1))
        _cell(pt.cell(ridx, 0), pr["date"], FD, bold=True, bg=GRAY, align="left"); ridx += 1
        for dft in pr["dfts"]:
            _cell(pt.cell(ridx, 0), "     " + dft["id"], FD, bg=WHITE, align="left")
            _cell(pt.cell(ridx, 1), dft["pct"], FD, bg=WHITE, align="right"); ridx += 1
    left_bottom = y + RH * n_plan

    # bloco esquerdo sem bordas (igual ao PNG)
    for _t in (vt, dt, pt):
        _sem_bordas(_t)

    CHART_X, CHART_Y, CHART_W = 6.15, 0.72, 6.95
    slide.shapes.add_picture(chart_png, Inches(CHART_X), Inches(CHART_Y), width=Inches(CHART_W))
    chart_bottom = CHART_Y + CHART_W / 2.08

    cols = d["detalhe"]["cols"]; drows = d["detalhe"]["rows"]
    dety = max(left_bottom, chart_bottom) + GAP
    LARG_DET = 13.333
    ws = [0.95, 1.45, 0.7, 0.9, 1.1, 3.9, 1.65, 1.35]
    ws += [(LARG_DET - sum(ws)) / len(cols[8:])] * len(cols[8:])
    COL_NOME_DET = 5                      # "Nome Defeito" dentro de row["cells"]
    f_base, f_nome, f_sec = 9.0, 8.0, 8.5

    # O nome do defeito vai inteiro e quebra em várias linhas, e o PowerPoint
    # estica a linha sozinho — era isso que jogava a tabela para fora do slide.
    # A altura de cada linha é estimada pelo tamanho do nome e as que não
    # couberem saem, preservando os defeitos de maior peso, que vêm primeiro.
    def _altura_det(nome):
        por_linha = max(int((ws[COL_NOME_DET] - 0.06) * 21.0), 10)
        linhas = max(1, -(-len(str(nome)) // por_linha))
        return max(0.24, linhas * f_nome * 1.25 / 72 + 0.06)

    disponivel = 7.5 - dety - 0.08
    rh_cab = 0.32
    usado_det = rh_cab
    linhas_ok = []
    for row in drows:
        h = rh_cab if row["kind"] == "sucesso" else _altura_det(row["cells"][COL_NOME_DET])
        if usado_det + h > disponivel:
            break
        usado_det += h
        linhas_ok.append((row, h))
    fora_det = len(drows) - len(linhas_ok)
    drows = [r for r, _ in linhas_ok]
    n = 1 + len(drows)

    dt2 = add_tbl(n, len(cols), 0.0, dety, LARG_DET, usado_det)
    _widths(dt2, ws)
    dt2.rows[0].height = Inches(rh_cab)
    for i, (_, h) in enumerate(linhas_ok):
        dt2.rows[1 + i].height = Inches(h)
    if fora_det:
        print(f"  [slide] {d['title'][:28]}: {fora_det} defeito(s) fora do slide, "
              f"descartados por falta de espaço")
    for ci, col in enumerate(cols):
        _cell(dt2.cell(0, ci), col, f_base, bold=True, fg=WHITE, bg=RED)
    for ri, row in enumerate(drows):
        bg = WHITE if row["kind"] == "sucesso" else (WHITE if (ri - 1) % 2 == 0 else GRAY2)
        for ci, val in enumerate(row["cells"]):
            al = "left" if ci == 5 else "center"
            fs = f_nome if ci == 5 else (f_sec if ci in (1, 6, 7) else f_base)
            bold = (ci == 0 and (row["kind"] == "sucesso" or val == "Falha"))
            _cell(dt2.cell(ri + 1, ci), val, fs, bold=bold, bg=bg, align=al)
    _sem_bordas(dt2)


# ── Slide "Top Ofensores": Top N defeitos/US de todas as jornadas juntas ────
JORNADA_LABEL_CURTO = {"Base Móvel": "Base"}


def _montar_slide_top_ofensores(prs, ns_consolidado, ns_por_jornada, top_n=15):
    """
    Adiciona um slide com uma única tabela "Top Ofensores": os `top_n` defeitos/US
    com maior % no mês mais recente, juntando todas as jornadas brutas em
    `ns_por_jornada` (cada uma vira uma linha "Jornada"), ordenados decrescente.
    Reaproveita o mesmo layout/estilo da tabela detalhada de `_montar_slide`.
    """
    import pandas as pd

    meses_tab = ns_consolidado["meses_tab"]
    MESES_LABEL = ns_consolidado["MESES_LABEL"]
    labels_tab = [MESES_LABEL[m][:3].capitalize() for m in meses_tab]
    resumo = ns_consolidado["resumo"]; pct_m_consolidado = ns_consolidado["pct_m"]

    candidatos = []
    for jornada, ns in ns_por_jornada.items():
        pivot2 = ns["pivot2"]; pct_m = ns["pct_m"]
        fmt_id = ns["fmt_id"]; fmt_ms = ns["fmt_milestone"]
        label_jornada = JORNADA_LABEL_CURTO.get(jornada, jornada)
        for _, row in pivot2.iterrows():
            tipo = str(row["DFT_Type"]).strip() if pd.notna(row["DFT_Type"]) else ""
            tipo_fmt = "US" if tipo == "User Story" else "Defeito"
            id_str = fmt_id({"DefectKey": row["DefectKey"],
                              "DFT_BugfixMilestone": row["DFT_BugfixMilestone"]})
            ms_fmt = fmt_ms(row["DFT_BugfixMilestone"])
            nome = str(row["DFT_Name"]) if pd.notna(row["DFT_Name"]) else ""
            team = str(row["DFT_Team"])[:35] if pd.notna(row["DFT_Team"]) else ""
            phase = str(row["DFT_Phase"])[:25] if pd.notna(row["DFT_Phase"]) else ""
            valores = [pct_m(int(row.loc[m]) if m in row.index else 0, m) for m in meses_tab]
            cells = [label_jornada, tipo_fmt, id_str, ms_fmt, nome, team, phase] + \
                    [f"{v:.2f}%".replace(".", ",") for v in valores]
            candidatos.append((valores[-1], cells))   # ordena pelo mês mais recente

    candidatos.sort(key=lambda c: c[0], reverse=True)
    top_linhas = [cells for _, cells in candidatos[:top_n]]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.55))
    tb.fill.solid(); tb.fill.fore_color.rgb = RED; tb.line.fill.background()
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Consolidado - Top Ofensores - Defeitos com Squad"
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

    def add_tbl(nr, nc, x, y, w, h):
        gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
        _no_style(gf.table); return gf.table

    COLS = ["Classificação", "Agrupamento", "Jornada", "Tipo", "Defeito/US",
            "Data Prod\nnão resolvidos", "Nome Defeito", "Time", "Status\nDefeito/US"] + labels_tab
    LARGURA = 13.333                       # ocupa o slide inteiro
    ws = [0.7, 1.6, 0.8, 0.6, 0.85, 1.0, 3.75, 1.45, 1.25]
    ws += [(LARGURA - sum(ws)) / len(labels_tab)] * len(labels_tab)
    COL_NOME = 6                           # índice de "Nome Defeito" em COLS

    y = 0.62
    FS_BASE, FS_SEC, FS_NOME = 10.0, 9.0, 8.0
    ALTURA_UTIL = 7.5 - y - 0.08           # até a borda de baixo do slide

    # O nome do defeito não é truncado, então a célula quebra em várias linhas e
    # o PowerPoint estica a linha sozinho — é o que jogava a tabela para fora do
    # slide. Aqui a altura de cada linha é estimada pelo tamanho do nome e as
    # que não couberem são descartadas, mantendo os defeitos de maior peso.
    def _altura(nome):
        # ~2,1 caracteres por 0,01" de largura na fonte 8pt
        por_linha = max(int((ws[COL_NOME] - 0.06) * 21.0), 10)
        linhas = max(1, -(-len(str(nome)) // por_linha))
        return max(0.30, linhas * FS_NOME * 1.25 / 72 + 0.06)

    rh_cab = 0.42                          # cabeçalho em duas linhas + Sucesso
    usado = rh_cab * 2
    cabem = []
    for cells in top_linhas:
        h = _altura(cells[4])              # 4 = Nome Defeito dentro de `cells`
        if usado + h > ALTURA_UTIL:
            break
        usado += h
        cabem.append((cells, h))
    cortados = len(top_linhas) - len(cabem)
    top_linhas = [c for c, _ in cabem]

    n_data = 1 + len(top_linhas)           # linha Sucesso + linhas que couberam
    n_total = 1 + n_data                   # + cabeçalho

    tbl = add_tbl(n_total, len(COLS), 0.0, y, LARGURA, usado)
    _widths(tbl, ws)
    tbl.rows[0].height = Inches(rh_cab)
    tbl.rows[1].height = Inches(rh_cab)
    for i, (_, h) in enumerate(cabem):
        tbl.rows[2 + i].height = Inches(h)

    for ci, col in enumerate(COLS):
        _cell(tbl.cell(0, ci), col, FS_BASE, bold=True, fg=WHITE, bg=RED)

    # ── Linha Sucesso (taxa consolidada, todas as jornadas somadas) ────────
    _cell(tbl.cell(1, 0), "Sucesso", FS_BASE, bg=WHITE)
    for ci in range(1, len(COLS) - len(labels_tab)):
        _cell(tbl.cell(1, ci), "", FS_BASE, bg=WHITE)
    for li, m in enumerate(meses_tab):
        suc_pct = pct_m_consolidado(resumo.loc[m, "Sucessos"] if m in resumo.index else 0, m)
        _cell(tbl.cell(1, len(COLS) - len(labels_tab) + li), f"{suc_pct:.2f}%".replace(".", ","), FS_BASE, bg=WHITE)

    # ── Linhas do Top N ──────────────────────────────────────────────────
    for ri, cells in enumerate(top_linhas):
        row_idx = 2 + ri
        bg = WHITE if ri % 2 == 0 else GRAY2
        _cell(tbl.cell(row_idx, 0), "Falha" if ri == 0 else "", FS_BASE, bg=bg, bold=(ri == 0))
        _cell(tbl.cell(row_idx, 1), "Em Tratamento/Avaliação\npela Squad" if ri == 0 else "", FS_SEC, bg=bg)
        for ci, val in enumerate(cells):
            al = "left" if ci == 4 else "center"       # 4 = Nome Defeito
            fs = FS_NOME if ci == 4 else (FS_SEC if ci in (0, 5, 6) else FS_BASE)
            _cell(tbl.cell(row_idx, 2 + ci), val, fs, bg=bg, align=al)

    _sem_bordas(tbl)
    if cortados:
        print(f"  [slide] Top Ofensores: {cortados} defeito(s) fora do slide, "
              f"descartados por falta de espaço")
