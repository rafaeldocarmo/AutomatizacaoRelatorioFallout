# -*- coding: utf-8 -*-
"""Relatório completo em PowerPoint editável (tabelas nativas + gráficos como
imagem): slide de 3 colunas, slide "Top Ofensores" e um slide por jornada.

Uso como módulo:  from gerar_slide3col import gerar_pptx_completo
                  buf = gerar_pptx_completo(base_dir)
"""
import os, io, tempfile, contextlib
import matplotlib
matplotlib.use("Agg")
import matplotlib.ticker as mticker
from matplotlib.backends.backend_agg import FigureCanvasAgg
from matplotlib.figure import Figure
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor

import fallout_core
import pipeline
import gerar_pptx
from gerar_pptx import (RED, WHITE, GRAY, FG,
                        _cell, _widths, _rowh, _no_style, _sem_bordas)

COLUNAS = [("Consolidado",       "Consolidado"),
           ("Prospect",          "Prospect PF"),
           ("Base + Cross Sell", "Base (Móvel) + Cross Sell")]

# Cores do gráfico (matplotlib usa string hex; WHITE/RED do pptx são RGBColor)
COR_BLUE  = "#1565C0"
COR_GRAYL = "#B0B0B0"
COR_BRANCO = "#FFFFFF"
COR_SEP   = "#DCDCDC"          # separador entre colunas (cinza esbranquiçado)
RGB_SEP   = RGBColor(0xDC, 0xDC, 0xDC)


# ── Extração ────────────────────────────────────────────────────────────────
def _extrair_col(jornada, base_dir):
    """Roda pipeline.gerar_relatorio() para a jornada e devolve os blocos da coluna."""
    with contextlib.redirect_stdout(io.StringIO()):
        ns = pipeline.gerar_relatorio(base_dir, jornada)

    resumo = ns["resumo"]; ML = ns["MESES_LABEL"]
    todos  = sorted(resumo.index.tolist())
    meses  = todos[-4:]
    fi = lambda v: f"{int(v):,}".replace(",", ".")
    fp = lambda v: f"{v:.2f}%".replace(".", ",")

    volume = {"months": [ML[m] for m in meses], "rows": [
        ("Vendas com Sucesso",      [fi(resumo.loc[m, "Sucessos"]) for m in meses]),
        ("Falha (Análise Técnica)", [fi(resumo.loc[m, "Falhas"])   for m in meses]),
        ("Fallout Rate",            [fp(resumo.loc[m, "Pct"])      for m in meses]),
    ]}

    # "Em Tratamento" + filhos no topo; demais por valor desc., ocultando zeros
    itens  = ns["dist_items"]                    # (label, valor, subitem, negrito)
    cabeca = [i for i in itens if not i[2]][:1]
    filhos = [i for i in itens if i[2]]
    resto  = sorted([i for i in itens if not i[2]][1:], key=lambda r: -r[1])
    resto  = [r for r in resto if r[1] > 0]
    dist = {"pct": ns["fallout_pct"],
            "rows": [(l, fp(v), s, b) for (l, v, s, b) in cabeca + filhos + resto]}

    red = ns["reducao"]
    plan = {"pct": red["Pct_plena"].sum(), "rows": [
        {"date": r["MilestoneDate"].strftime("%d/%m/%Y"),
         "dfts": [d.strip() for d in r["DFTs"].split("\n")],
         "pct":  fp(r["Pct"])} for _, r in red.iterrows()]}

    chart = {"labels":  [ML[m] for m in todos],
             "vals":    [resumo.loc[m, "Pct"] for m in todos],
             "plabels": ns["proj_labels"], "pvals": list(ns["proj_vals"])}

    return {"volume": volume, "dist": dist, "plan": plan, "chart": chart,
            "corte": ns["hoje"].tz_convert("America/Sao_Paulo"), "ns": ns}


def _chart_png(c, destino, larg_in, alt_in):
    """Gráfico compacto da coluna (real + projeção + meta)."""
    fig = Figure(figsize=(larg_in, alt_in))
    FigureCanvasAgg(fig)
    fig.patch.set_facecolor(COR_BRANCO)
    ax = fig.add_axes([0.13, 0.26, 0.84, 0.66])
    labs = c["labels"] + c["plabels"][1:]

    ax.plot(range(len(c["labels"])), c["vals"], color="black", lw=1.3,
            marker="o", ms=3.4, zorder=3)
    if len(c["plabels"]) > 1:
        xs = list(range(len(c["labels"]) - 1, len(labs)))
        ax.plot(xs, c["pvals"], color=COR_GRAYL, lw=1.6, ls="--",
                marker="o", ms=4.2, mfc=COR_BRANCO, zorder=2)
        ax.annotate(f"{c['pvals'][-1]:.2f}%".replace(".", ","),
                    xy=(xs[-1], c["pvals"][-1]), xytext=(4, 0), textcoords="offset points",
                    ha="left", va="center", fontsize=7.5, fontweight="bold")
    for i, v in enumerate(c["vals"]):
        ax.annotate(f"{v:.2f}%".replace(".", ","), xy=(i, v), xytext=(0, 6),
                    textcoords="offset points", ha="center", fontsize=6.2, rotation=90)

    ax.axhline(y=1, color=COR_BLUE, lw=1.2)
    ax.set_ylim(0, max(c["vals"] + list(c["pvals"])) * 1.55)
    ax.set_xlim(-0.6, len(labs) + 0.7)
    ax.set_xticks(range(len(labs)))
    ax.set_xticklabels(labs, rotation=45, ha="right", fontsize=5.2)
    ax.tick_params(axis="y", labelsize=6, length=2)
    _dec = 0 if (ax.get_ylim()[1] - ax.get_ylim()[0]) >= 5 else 1
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _, d=_dec: f"{v:.{d}f}%".replace(".", ",")))
    ax.grid(axis="y", ls="--", alpha=0.25)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)

    fig.savefig(destino, dpi=200, facecolor=COR_BRANCO)


def extrair_todas(base_dir="."):
    """Roda o pipeline das 3 visões uma única vez (reaproveitável nos 2 formatos)."""
    return {j: _extrair_col(j, base_dir) for j, _ in COLUNAS}


# ── Montagem: PowerPoint ────────────────────────────────────────────────────
def _montar_slide3col(prs, dados):
    """Adiciona o slide de 3 colunas à apresentação `prs` já existente."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Faixa de título
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.52))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    p = bar.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Realização e Projeção de Fallout"
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

    corte = list(dados.values())[0]["corte"]
    tb = slide.shapes.add_textbox(Inches(0.12), Inches(0.54), Inches(2.2), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = f"Data de Corte: {corte.strftime('%d/%b')}"
    r.font.size = Pt(9); r.font.bold = True; r.font.italic = True
    r.font.color.rgb = FG; r.font.name = "Calibri"

    M, GAP = 0.15, 0.22
    CW = (13.333 - 2*M - 2*GAP) / 3
    XS = [M + i*(CW + GAP) for i in range(3)]

    Y_CHART, H_CHART = 0.88, 1.42
    Y_TAB, Y_FUNDO, GAP_BL = 2.42, 7.35, 0.06

    def n_linhas(ci, d):
        n = (2 + len(d["volume"]["rows"])) if ci > 0 else 0
        n += 1 + len(d["dist"]["rows"])
        n += 1 + sum(len(r["dfts"]) for r in d["plan"]["rows"])
        return n

    max_lin = max(n_linhas(ci, dados[j]) for ci, (j, _) in enumerate(COLUNAS))
    RH = min(0.20, (Y_FUNDO - Y_TAB - 2*GAP_BL) / max_lin)
    # Fontes conforme a altura de linha (RH em pt = RH*72; a linha ocupa ~1.2*fonte)
    if RH >= 0.16:   FD, FT = 8.0, 9.0     # dados, títulos
    elif RH >= 0.13: FD, FT = 7.0, 8.0
    else:            FD, FT = 6.0, 7.0

    tmp = tempfile.mkdtemp(prefix="slide3col_")

    def add_tbl(nr, nc, x, y, w, h):
        gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
        _no_style(gf.table); return gf.table

    # Separadores pontilhados entre as colunas
    for ci in (1, 2):
        xs = XS[ci] - GAP/2
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(xs), Inches(0.60),
                                        Inches(xs), Inches(Y_FUNDO))
        cn.line.color.rgb = RGB_SEP
        cn.line.width = Pt(1)
        cn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    for ci, (jornada, titulo) in enumerate(COLUNAS):
        d = dados[jornada]; x0 = XS[ci]

        # Título da coluna
        tb = slide.shapes.add_textbox(Inches(x0), Inches(0.52), Inches(CW), Inches(0.32))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = titulo
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = FG; r.font.name = "Calibri"

        # Gráfico
        png = os.path.join(tmp, f"c{ci}.png")
        _chart_png(d["chart"], png, CW, H_CHART)
        slide.shapes.add_picture(png, Inches(x0), Inches(Y_CHART), width=Inches(CW))

        y = Y_TAB

        # ── Volume de Pedidos (colunas 2 e 3) ──
        if ci > 0:
            v = d["volume"]; n = len(v["months"])
            nr = 2 + len(v["rows"])
            t = add_tbl(nr, 1 + n, x0, y, CW, RH*nr)
            lw_ = CW * 0.40
            _widths(t, [lw_] + [(CW - lw_)/n]*n); _rowh(t, RH)
            t.cell(0, 0).merge(t.cell(0, n))
            _cell(t.cell(0, 0), "Volume de Pedidos", FT, bold=True, fg=WHITE, bg=RED)
            _cell(t.cell(1, 0), "", FD, bg=RED)
            for i, m in enumerate(v["months"]):
                _cell(t.cell(1, 1 + i), m, FD, bold=True, fg=WHITE, bg=RED)
            for ri, (lbl, vals) in enumerate(v["rows"]):
                bg = WHITE if ri % 2 == 0 else GRAY
                _cell(t.cell(2 + ri, 0), lbl, FD, bg=bg, align="left")
                for i, val in enumerate(vals):
                    _cell(t.cell(2 + ri, 1 + i), val, FD, bg=bg)
            _sem_bordas(t)
            y += RH*nr + GAP_BL

        # ── Distribuição Fallout ──
        rows = d["dist"]["rows"]; nr = 1 + len(rows)
        t = add_tbl(nr, 2, x0, y, CW, RH*nr)
        VW = CW * 0.26
        _widths(t, [CW - VW, VW]); _rowh(t, RH)
        t.cell(0, 0).merge(t.cell(0, 1))
        _cell(t.cell(0, 0), f"Distribuição Fallout ({d['dist']['pct']:.2f}%)".replace(".", ","),
              FT, bold=True, fg=WHITE, bg=RED)
        for ri, (lbl, val, sub, _b) in enumerate(rows):
            bg = WHITE if ri % 2 == 0 else GRAY
            # toda a Distribuição em negrito
            _cell(t.cell(1 + ri, 0), lbl, FD, bold=True, bg=bg,
                  align="right" if sub else "left")
            _cell(t.cell(1 + ri, 1), val, FD, bold=True, bg=bg, align="right")
        _sem_bordas(t)
        y += RH*nr + GAP_BL

        # ── Planejamento Redução ──
        prows = d["plan"]["rows"]
        nr = 1 + len(prows)
        alturas = [RH] + [RH*len(r["dfts"]) for r in prows]
        t = add_tbl(nr, 3, x0, y, CW, sum(alturas))
        DW, PW = CW*0.30, CW*0.22
        _widths(t, [DW, CW - DW - PW, PW])
        for i, h in enumerate(alturas):
            t.rows[i].height = Inches(h)
        t.cell(0, 0).merge(t.cell(0, 2))
        _cell(t.cell(0, 0), f"Planejamento Redução ({d['plan']['pct']:.2f}%)".replace(".", ","),
              FT, bold=True, fg=WHITE, bg=RED)
        for ri, r in enumerate(prows):
            bg = WHITE if ri % 2 == 0 else GRAY
            _cell(t.cell(1 + ri, 0), r["date"], FD, bg=bg)
            _cell(t.cell(1 + ri, 1), "\n".join(r["dfts"]), FD, bold=True, bg=bg)
            _cell(t.cell(1 + ri, 2), r["pct"], FD, bold=True, bg=bg, align="right")
        _sem_bordas(t)


def gerar_pptx_completo(base_dir=".", dados_3col=None):
    """
    PPTX único:
      1. Slide de 3 colunas (Consolidado | Prospect | Base + Cross Sell)
      2. Slide "Top Ofensores" — Top 15 defeitos/US de todas as jornadas juntas
      3. Um slide no estilo do "print" para cada jornada bruta encontrada em
         extrações/ (ex.: Base Móvel, Cross Sell, Prospect — sem os combinados)
    Devolve os bytes (BytesIO).
    """
    dados_3col = dados_3col or extrair_todas(base_dir)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    _montar_slide3col(prs, dados_3col)

    jornadas_brutas = fallout_core.jornadas_disponiveis(base_dir)
    ns_por_jornada = {j: pipeline.gerar_relatorio(base_dir, j) for j in jornadas_brutas}

    # O Top Ofensores é a visão consolidada, então lista as mesmas jornadas que
    # compõem o Consolidado. As de formato RPA (PME) entram apenas com o slide
    # individual delas, mais abaixo.
    consolidaveis = set(fallout_core.jornadas_consolidaveis(base_dir))
    ns_consolidado = (dados_3col.get("Consolidado") or {}).get("ns")
    if ns_consolidado is None:
        ns_consolidado = pipeline.gerar_relatorio(base_dir, "Consolidado")
    gerar_pptx._montar_slide_top_ofensores(
        prs, ns_consolidado,
        {j: ns for j, ns in ns_por_jornada.items() if j in consolidaveis},
    )

    for jornada, ns in ns_por_jornada.items():
        data, chart_png = gerar_pptx._montar_dados(jornada, ns)
        gerar_pptx._montar_slide(prs, data, chart_png)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf
